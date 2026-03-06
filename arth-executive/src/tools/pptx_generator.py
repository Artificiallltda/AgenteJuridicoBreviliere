import os
import json
import uuid
import logging
import base64
import httpx
import traceback
from langchain_core.tools import tool
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from src.config import settings

logger = logging.getLogger(__name__)

def _apply_dark_theme(slide):
    """Aplica fundo escuro moderno (Dark Mode) em um slide específico."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(33, 37, 41)

def _style_run(run, size_pt, bold=False, rgb=RGBColor(255, 255, 255)):
    """Formata a fonte de um fragmento de texto."""
    run.font.name = "Segoe UI"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = rgb

@tool
def generate_pptx(slides_content_json: str) -> str:
    """
    Cria uma apresentação de Slides (.pptx).
    A entrada `slides_content_json` deve ser uma string com um JSON válido contendo a estrutura dos slides.
    Formato esperado:
    {
      "presentation_title": "Título Geral da Apresentação",
      "slides": [
         {
           "title": "Título do Slide 1",
           "bullets": ["Tópico 1", "Tópico 2", "Tópico 3"],
           "image_prompt": "Uma foto realista e futurista de IA (Opcional)"
         }
      ]
    }
    """
    try:
        output_dir = os.path.join(os.getcwd(), "data", "outputs")
        os.makedirs(output_dir, exist_ok=True)
        
        filename = f"Apresentacao_{uuid.uuid4().hex[:6]}.pptx"
        filepath = os.path.join(output_dir, filename)

        try:
            content = json.loads(slides_content_json)
        except Exception as e:
            return f"Erro de JSON: Não foi possível realizar o parse da string fornecida. Garanta que mandou um JSON perfeitamente válido. Falha: {e}"

        prs = Presentation()
        # Garante proporção Widescreen 16:9 (Padrão corporativo)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        presentation_title = content.get("presentation_title", "Apresentação Gerada")
        slides = content.get("slides", [])

        # ==========================================
        # CAPA
        # ==========================================
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        _apply_dark_theme(slide)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        # Centraliza a capa estritamente
        title_shape.left = Inches(1.0)
        title_shape.top = Inches(2.5)
        title_shape.width = Inches(11.33)
        title_shape.height = Inches(1.5)
        
        title_shape.text = presentation_title
        if title_shape.text_frame and title_shape.text_frame.paragraphs and title_shape.text_frame.paragraphs[0].runs:
            _style_run(title_shape.text_frame.paragraphs[0].runs[0], size_pt=54, bold=True, rgb=RGBColor(255, 255, 255))

        subtitle_shape.left = Inches(1.0)
        subtitle_shape.top = Inches(4.5)
        subtitle_shape.width = Inches(11.33)
        subtitle_shape.height = Inches(1.0)
        subtitle_shape.text = "Gerado por Arth Executive AI"
        if subtitle_shape.text_frame and subtitle_shape.text_frame.paragraphs and subtitle_shape.text_frame.paragraphs[0].runs:
            _style_run(subtitle_shape.text_frame.paragraphs[0].runs[0], size_pt=24, bold=False, rgb=RGBColor(180, 180, 180))

        # ==========================================
        # SLIDES DE CONTEÚDO
        # ==========================================
        bullet_slide_layout = prs.slide_layouts[1]
        from src.tools.image_generator import generate_image
        import re

        for slide_data in slides:
            slide_title = slide_data.get("title", "")
            bullets = slide_data.get("bullets", [])
            img_prompt = slide_data.get("image_prompt", "")

            slide = prs.slides.add_slide(bullet_slide_layout)
            _apply_dark_theme(slide)
            
            # 1. Ajusta o Título rigorosamente no TOPO
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.left = Inches(0.5)
                title_shape.top = Inches(0.4)
                title_shape.width = Inches(12.33)
                title_shape.height = Inches(1.0)
                title_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
                title_shape.text = slide_title
                if title_shape.text_frame and title_shape.text_frame.paragraphs and title_shape.text_frame.paragraphs[0].runs:
                    _style_run(title_shape.text_frame.paragraphs[0].runs[0], size_pt=40, bold=True, rgb=RGBColor(255, 255, 255))

            # 2. Ajusta a Caixa de Texto dos Bullets Abaixo do Título
            body_shape = slide.placeholders[1]
            body_shape.left = Inches(0.5)
            body_shape.top = Inches(1.6)
            body_shape.width = Inches(12.33)
            body_shape.height = Inches(5.2)
            body_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            
            # Se houver imagem, baixa e realoca o design para Side-by-Side Horizontal
            img_path = None
            if img_prompt and settings.OPENAI_API_KEY:
                logger.info(f"[PPTX] Gerando imagem para slide via gpt-image-1.5: {img_prompt[:30]}...")
                try:
                    # Invocamos a tool já robusta (que lida com os webhooks do AIOS)
                    # Ela salvará na pasta outputs e retornará uma string com o <SEND_FILE:nome.png>
                    result_str = generate_image.invoke({"prompt": img_prompt})
                    
                    # Extrai o nome do arquivo da string de resposta
                    match = re.search(r'<SEND_FILE:([^>]+)>', result_str)
                    if match:
                        img_filename = match.group(1)
                        img_path = os.path.join(output_dir, img_filename)
                        
                        if not os.path.exists(img_path):
                            logger.warning(f"O gerador retornou {img_filename}, mas o arquivo não foi encontrado.")
                            img_path = None
                    else:
                        logger.warning(f"Falha ao extrair a tag de arquivo da resposta do DALL-E: {result_str}")
                except Exception as img_err:
                    logger.warning(f"Exceção ao gerar imagem: {img_err}")

            if img_path:
                # Diminui a caixa de texto pra METADE pro texto não atropelar a imagem
                body_shape.width = Inches(6.0)
                
                # Renderiza a imagem perfeitamente enquadrada na COLUNA DIREITA
                pic_left = Inches(6.8)
                pic_top = Inches(1.6)
                pic_height = Inches(5.2)
                
                try:
                    # Injeta a foto. Manter o height e auto-ajustar a width mantém a proporção 1:1 original do DALL-E do lado direito
                    slide.shapes.add_picture(img_path, pic_left, pic_top, height=pic_height)
                except Exception as e:
                    logger.error(f"Não foi possível aplicar a imagem ao slide, ignorando... {e}")

            # Preenche os bullets
            tf = body_shape.text_frame
            for i, bullet_text in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                    p.text = str(bullet_text)
                else:
                    p = tf.add_paragraph()
                    p.text = str(bullet_text)
                    p.level = 0
                
                # Margem entre linhas
                p.space_after = Pt(14)
                
                if p.runs:
                    _style_run(p.runs[0], size_pt=24, bold=False, rgb=RGBColor(220, 220, 220))

        prs.save(filepath)
        
        return f"A apresentação foi montada em disco. Use SEMPRE a tag <SEND_FILE:{filename}> na sua resposta ao usuário para enviar o arquivo."

    except Exception as e:
        logger.error(f"Falha ao gerar o PPTX: {e}")
        return f"Ocorreu um erro interno ao compilar os slides da apresentação. Detalhes: {e}\n{traceback.format_exc()}"
